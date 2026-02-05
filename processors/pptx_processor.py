"""
PPTX presentation processor for PowerPoint files.
"""

import io
import logging
from pathlib import Path
from typing import Generator

from pptx import Presentation
from pptx.shapes.base import BaseShape
from pptx.shapes.group import GroupShape
from pptx.shapes.picture import Picture
from pptx.util import Inches, Pt

from .base import BaseFileProcessor, ContentChunk, ContentType
from settings import settings

logger = logging.getLogger("OfficeTranslator.PptxProcessor")


class PptxProcessor(BaseFileProcessor):
    """Processor for Microsoft PowerPoint (.pptx) presentations."""
    
    SUPPORTED_EXTENSIONS = ["pptx"]
    
    def __init__(self):
        super().__init__()
        self._document = None
        self._shape_map: dict[str, tuple] = {}  # id -> (slide_idx, shape reference)
        self._image_map: dict[str, tuple] = {}  # id -> (slide_idx, picture shape)
    
    def load(self, file_path: str | Path) -> None:
        """Load a PPTX file."""
        path = Path(file_path)
        
        if not path.exists():
            raise FileNotFoundError(f"File not found: {path}")
        
        if path.suffix.lower() != ".pptx":
            raise ValueError(f"Not a PPTX file: {path}")
        
        try:
            self._document = Presentation(path)
            self._file_path = path
            self._shape_map.clear()
            self._image_map.clear()
            logger.info(f"Loaded PPTX: {path.name} ({len(self._document.slides)} slides)")
        except Exception as e:
            raise ValueError(f"Invalid PPTX file: {e}")
    
    def extract_content_generator(self) -> Generator[ContentChunk, None, None]:
        """Extract content from all slides."""
        self._validate_loaded()
        
        chunk_index = 0
        
        for slide_idx, slide in enumerate(self._document.slides):
            slide_num = slide_idx + 1
            
            # Process shapes recursively (handles groups)
            for shape_idx, shape in enumerate(slide.shapes):
                yield from self._process_shape(
                    shape, slide_idx, shape_idx, slide_num, chunk_index
                )
        
        logger.info(f"Extracted {len(self._shape_map) + len(self._image_map)} content chunks from PPTX")
    
    def _process_shape(
        self,
        shape: BaseShape,
        slide_idx: int,
        shape_idx: int,
        slide_num: int,
        chunk_index: int,
        parent_prefix: str = "",
    ) -> Generator[ContentChunk, None, None]:
        """Process a single shape, recursively handling groups."""
        
        shape_prefix = f"{parent_prefix}s{shape_idx}" if parent_prefix else f"s{shape_idx}"
        
        # Handle grouped shapes
        if isinstance(shape, GroupShape):
            for sub_idx, sub_shape in enumerate(shape.shapes):
                yield from self._process_shape(
                    sub_shape, slide_idx, sub_idx, slide_num,
                    chunk_index, f"{shape_prefix}_g"
                )
            return
        
        # Handle pictures/images
        if isinstance(shape, Picture):
            try:
                image_data = shape.image.blob
                img_id = f"slide{slide_idx}_{shape_prefix}_img"
                self._image_map[img_id] = (slide_idx, shape)
                
                yield ContentChunk(
                    id=img_id,
                    content_type=ContentType.IMAGE,
                    image_data=image_data,
                    location=f"Slide {slide_num}, Image",
                    metadata={
                        "slide": slide_idx,
                        "content_type": shape.image.content_type,
                    }
                )
            except Exception as e:
                logger.debug(f"Failed to extract image: {e}")
            return
        
        # Handle text frames
        if shape.has_text_frame:
            text_frame = shape.text_frame
            
            for para_idx, paragraph in enumerate(text_frame.paragraphs):
                para_text = paragraph.text.strip()
                if para_text:
                    chunk_id = f"slide{slide_idx}_{shape_prefix}_p{para_idx}"
                    self._shape_map[chunk_id] = (slide_idx, paragraph)
                    
                    yield ContentChunk(
                        id=chunk_id,
                        content_type=ContentType.TEXT,
                        text=para_text,
                        location=f"Slide {slide_num}",
                        metadata={
                            "slide": slide_idx,
                            "shape": shape_idx,
                            "paragraph": para_idx,
                            "has_formatting": self._has_formatting(paragraph),
                        }
                    )
        
        # Handle tables
        if shape.has_table:
            table = shape.table
            for row_idx, row in enumerate(table.rows):
                for cell_idx, cell in enumerate(row.cells):
                    cell_text = cell.text.strip()
                    if cell_text:
                        chunk_id = f"slide{slide_idx}_{shape_prefix}_t_r{row_idx}_c{cell_idx}"
                        # Store reference to first paragraph in cell
                        self._shape_map[chunk_id] = (slide_idx, cell.text_frame.paragraphs[0])
                        
                        yield ContentChunk(
                            id=chunk_id,
                            content_type=ContentType.TABLE_CELL,
                            text=cell_text,
                            location=f"Slide {slide_num}, Table",
                            metadata={
                                "slide": slide_idx,
                                "row": row_idx,
                                "col": cell_idx,
                            }
                        )
    
    def apply_translation(
        self,
        chunk_id: str,
        translated_content: str | bytes,
    ) -> None:
        """Apply translation to a text shape or image."""
        self._validate_loaded()
        
        if chunk_id in self._image_map:
            # Handle image replacement
            if isinstance(translated_content, bytes):
                slide_idx, picture_shape = self._image_map[chunk_id]
                self._replace_picture(picture_shape, translated_content)
            return
        
        if chunk_id not in self._shape_map:
            logger.warning(f"Chunk not found: {chunk_id}")
            return
        
        slide_idx, paragraph = self._shape_map[chunk_id]
        
        if settings.preserve_formatting and self._has_formatting(paragraph):
            self._apply_with_formatting(paragraph, str(translated_content))
        else:
            # Simple replacement
            self._simple_replace(paragraph, str(translated_content))
        
        logger.debug(f"Applied translation to {chunk_id}")
    
    def save(self, output_path: str | Path) -> None:
        """Save the modified presentation."""
        self._validate_loaded()
        
        path = Path(output_path)
        self._document.save(path)
        logger.info(f"Saved PPTX: {path.name}")
    
    def _has_formatting(self, paragraph) -> bool:
        """Check if paragraph has inline formatting."""
        for run in paragraph.runs:
            font = run.font
            if font.bold or font.italic or font.underline:
                return True
        return False
    
    def _simple_replace(self, paragraph, text: str) -> None:
        """Replace paragraph text, preserving first run's basic formatting."""
        if paragraph.runs:
            # Keep font settings from first run
            first_run = paragraph.runs[0]
            font_name = first_run.font.name
            font_size = first_run.font.size
            font_bold = first_run.font.bold
            font_italic = first_run.font.italic
            
            # Clear all runs
            for run in paragraph.runs:
                run.text = ""
            
            # Set text on first run
            paragraph.runs[0].text = text
            paragraph.runs[0].font.name = font_name
            if font_size:
                paragraph.runs[0].font.size = font_size
            paragraph.runs[0].font.bold = font_bold
            paragraph.runs[0].font.italic = font_italic
        else:
            paragraph.text = text
    
    def _apply_with_formatting(self, paragraph, text: str) -> None:
        """Apply translation while attempting to preserve formatting."""
        # For now, use simple replacement
        # Full implementation would parse formatting tags from translation
        self._simple_replace(paragraph, text)
    
    def _replace_picture(self, picture_shape: Picture, new_image_data: bytes) -> None:
        """
        Replace a picture shape's image.
        
        Note: This is a simplified implementation. Full implementation
        would need to handle maintaining the original size/position.
        """
        try:
            # Get original dimensions
            width = picture_shape.width
            height = picture_shape.height
            left = picture_shape.left
            top = picture_shape.top
            
            # The python-pptx library doesn't have a direct way to replace
            # the image blob, so we log this as a limitation
            logger.warning("Image replacement in PPTX requires additional implementation")
            
        except Exception as e:
            logger.error(f"Failed to replace picture: {e}")
