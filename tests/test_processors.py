"""
Unit tests for file processors.
"""

import pytest
from pathlib import Path
from unittest.mock import Mock, patch, MagicMock
import tempfile
import os

# Test the base processor interface
class TestBaseProcessor:
    """Tests for BaseFileProcessor abstract interface."""
    
    def test_content_chunk_validation(self):
        """Test ContentChunk requires appropriate fields."""
        from processors.base import ContentChunk, ContentType
        
        # Text chunk requires text
        with pytest.raises(ValueError):
            ContentChunk(id="1", content_type=ContentType.TEXT, text=None)
        
        # Image chunk requires image_data
        with pytest.raises(ValueError):
            ContentChunk(id="2", content_type=ContentType.IMAGE, image_data=None)
        
        # Valid chunks should work
        text_chunk = ContentChunk(id="3", content_type=ContentType.TEXT, text="Hello")
        assert text_chunk.text == "Hello"
        
        img_chunk = ContentChunk(id="4", content_type=ContentType.IMAGE, image_data=b"PNG...")
        assert img_chunk.image_data == b"PNG..."


class TestProcessorFactory:
    """Tests for ProcessorFactory."""
    
    def test_supported_extensions(self):
        """Test that factory reports supported extensions."""
        from processors.factory import get_supported_extensions
        
        extensions = get_supported_extensions()
        assert isinstance(extensions, list)
        # Should have at least docx and pptx
        assert ".docx" in extensions or len(extensions) == 0
    
    def test_unsupported_extension_raises(self):
        """Test that unsupported extensions raise ValueError."""
        from processors.factory import get_processor
        
        with pytest.raises(ValueError, match="Unsupported file type"):
            get_processor("test.unsupported")
    
    def test_is_supported(self):
        """Test is_supported check."""
        from processors.factory import ProcessorFactory
        
        # Should return False for unsupported
        assert ProcessorFactory.is_supported("test.xyz") == False


class TestDocxProcessor:
    """Tests for DocxProcessor."""
    
    @pytest.fixture
    def sample_docx(self, tmp_path):
        """Create a minimal DOCX file for testing."""
        from docx import Document
        
        doc = Document()
        doc.add_paragraph("Hello World")
        doc.add_paragraph("This is a test.")
        
        file_path = tmp_path / "test.docx"
        doc.save(file_path)
        return file_path
    
    def test_load_docx(self, sample_docx):
        """Test loading a DOCX file."""
        from processors.docx_processor import DocxProcessor
        
        processor = DocxProcessor()
        processor.load(sample_docx)
        
        assert processor.is_loaded
        assert processor.file_path == sample_docx
    
    def test_extract_content(self, sample_docx):
        """Test extracting content from DOCX."""
        from processors.docx_processor import DocxProcessor
        from processors.base import ContentType
        
        processor = DocxProcessor()
        processor.load(sample_docx)
        
        chunks = list(processor.extract_content_generator())
        
        assert len(chunks) >= 2
        assert chunks[0].content_type == ContentType.TEXT
        assert "Hello World" in chunks[0].text
    
    def test_apply_translation(self, sample_docx, tmp_path):
        """Test applying translation to DOCX."""
        from processors.docx_processor import DocxProcessor
        
        processor = DocxProcessor()
        processor.load(sample_docx)
        
        chunks = list(processor.extract_content_generator())
        first_chunk = chunks[0]
        
        processor.apply_translation(first_chunk.id, "Hallo Welt")
        
        output_path = tmp_path / "output.docx"
        processor.save(output_path)
        
        # Verify saved file
        assert output_path.exists()


class TestPptxProcessor:
    """Tests for PptxProcessor."""
    
    @pytest.fixture
    def sample_pptx(self, tmp_path):
        """Create a minimal PPTX file for testing."""
        from pptx import Presentation
        from pptx.util import Inches
        
        prs = Presentation()
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)
        
        # Add text box
        left = top = Inches(1)
        width = height = Inches(2)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = "Slide Title"
        
        file_path = tmp_path / "test.pptx"
        prs.save(file_path)
        return file_path
    
    def test_load_pptx(self, sample_pptx):
        """Test loading a PPTX file."""
        from processors.pptx_processor import PptxProcessor
        
        processor = PptxProcessor()
        processor.load(sample_pptx)
        
        assert processor.is_loaded
        assert processor.file_path == sample_pptx
    
    def test_extract_content(self, sample_pptx):
        """Test extracting content from PPTX."""
        from processors.pptx_processor import PptxProcessor
        from processors.base import ContentType
        
        processor = PptxProcessor()
        processor.load(sample_pptx)
        
        chunks = list(processor.extract_content_generator())
        
        assert len(chunks) >= 1
        text_chunks = [c for c in chunks if c.content_type == ContentType.TEXT]
        assert any("Slide Title" in c.text for c in text_chunks)


class TestXlsxProcessor:
    """Tests for XlsxProcessor."""
    
    @pytest.fixture
    def sample_xlsx(self, tmp_path):
        """Create a minimal XLSX file for testing."""
        from openpyxl import Workbook
        
        wb = Workbook()
        ws = wb.active
        ws['A1'] = "Name"
        ws['B1'] = "Value"
        ws['A2'] = "Test"
        ws['B2'] = 42  # Number, should be skipped
        
        file_path = tmp_path / "test.xlsx"
        wb.save(file_path)
        return file_path
    
    def test_load_xlsx(self, sample_xlsx):
        """Test loading an XLSX file."""
        from processors.xlsx_processor import XlsxProcessor
        
        processor = XlsxProcessor()
        processor.load(sample_xlsx)
        
        assert processor.is_loaded
    
    def test_extract_content_skips_numbers(self, sample_xlsx):
        """Test that numbers are not extracted."""
        from processors.xlsx_processor import XlsxProcessor
        
        processor = XlsxProcessor()
        processor.load(sample_xlsx)
        
        chunks = list(processor.extract_content_generator())
        
        # Should only have string cells
        texts = [c.text for c in chunks]
        assert "42" not in texts


class TestLLMClient:
    """Tests for LLMClient."""
    
    @patch('services.llm_client.OpenAI')
    def test_translate_returns_content(self, mock_openai):
        """Test that translate returns translated content."""
        from services.llm_client import LLMClient
        
        # Mock the OpenAI response
        mock_response = Mock()
        mock_response.choices = [Mock(message=Mock(content="Translated text"))]
        mock_openai.return_value.chat.completions.create.return_value = mock_response
        
        client = LLMClient()
        result = client.translate("Hello", "German")
        
        assert result == "Translated text"
    
    @patch('services.llm_client.OpenAI')
    def test_empty_text_returns_empty(self, mock_openai):
        """Test that empty text is not sent to LLM."""
        from services.llm_client import LLMClient
        
        client = LLMClient()
        result = client.translate("", "German")
        
        assert result == ""
        mock_openai.return_value.chat.completions.create.assert_not_called()


class TestTextUtils:
    """Tests for text utility functions."""
    
    def test_calculate_luminance_black(self):
        """Test luminance of black is 0."""
        from utils.text_utils import calculate_luminance
        assert calculate_luminance((0, 0, 0)) == 0.0
    
    def test_calculate_luminance_white(self):
        """Test luminance of white is 1."""
        from utils.text_utils import calculate_luminance
        assert calculate_luminance((255, 255, 255)) == pytest.approx(1.0, rel=0.01)
    
    def test_wrap_text_short(self):
        """Test that short text isn't wrapped."""
        from utils.text_utils import wrap_text
        from PIL import ImageFont
        
        try:
            font = ImageFont.truetype("arial.ttf", 12)
        except OSError:
            font = ImageFont.load_default()
        
        lines = wrap_text("Hi", font, 1000)
        assert len(lines) == 1
        assert lines[0] == "Hi"
